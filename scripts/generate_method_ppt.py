#!/usr/bin/env python3

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
RESULTS = ROOT / "results"
DOCS = ROOT / "docs"
OUTPUT_PPTX = DOCS / "method_section_20260323.pptx"
OUTPUT_PPTX_CN = DOCS / "方法部分_空中对接_20260323.pptx"

COLOR_BG = RGBColor(246, 248, 252)
COLOR_NAVY = RGBColor(15, 30, 60)
COLOR_BLUE = RGBColor(35, 94, 184)
COLOR_LIGHT_BLUE = RGBColor(218, 232, 252)
COLOR_GREEN = RGBColor(29, 125, 72)
COLOR_ORANGE = RGBColor(205, 121, 35)
COLOR_RED = RGBColor(180, 55, 42)
COLOR_TEXT = RGBColor(32, 37, 43)
COLOR_MUTED = RGBColor(96, 106, 120)
COLOR_WHITE = RGBColor(255, 255, 255)


def read_kv(path: Path) -> dict[str, str]:
    data: dict[str, str] = {}
    for line in path.read_text(encoding="utf-8", errors="ignore").splitlines():
        if "=" in line:
            key, value = line.split("=", 1)
            data[key.strip()] = value.strip()
    return data


def set_text_style(run, size: float, bold: bool = False, color: RGBColor = COLOR_TEXT,
                   font_name: str = "Microsoft YaHei") -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_textbox(slide, left, top, width, height, text="", size=20, bold=False, color=COLOR_TEXT,
                fill_color: RGBColor | None = None, line_color: RGBColor | None = None,
                align=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.TOP, margin=0.10):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = vertical
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_text_style(run, size=size, bold=bold, color=color)

    if fill_color is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()

    if line_color is not None:
        box.line.color.rgb = line_color
    else:
        box.line.fill.background()
    return box


def add_bullets(slide, left, top, width, height, bullets: Iterable[str], title: str | None = None,
                fill_color: RGBColor | None = None, line_color: RGBColor | None = None,
                title_color: RGBColor = COLOR_BLUE, bullet_size: float = 18.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)

    if fill_color is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()
    if line_color is not None:
        box.line.color.rgb = line_color
    else:
        box.line.fill.background()

    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        set_text_style(run, 16, bold=True, color=title_color)
    else:
        tf.clear()

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            set_text_style(p.runs[0], bullet_size, color=COLOR_TEXT)
    return box


def add_metric_card(slide, left, top, width, height, title: str, value: str,
                    accent: RGBColor = COLOR_BLUE, note: str | None = None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.10)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_text_style(r1, 12, bold=True, color=accent)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    set_text_style(r2, 22, bold=True, color=COLOR_TEXT)
    if note:
        p3 = tf.add_paragraph()
        r3 = p3.add_run()
        r3.text = note
        set_text_style(r3, 10, color=COLOR_MUTED)
    return shape


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(11.6), Inches(0.6), title, size=26, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.62), Inches(0.92), Inches(11.3), Inches(0.35), subtitle, size=10.5, color=COLOR_MUTED)


def add_footer(slide, text: str) -> None:
    add_textbox(slide, Inches(0.6), Inches(6.95), Inches(11.8), Inches(0.2), text, size=9, color=COLOR_MUTED)


def add_image(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    mock_summary = read_kv(RESULTS / "20260318_120407_mock_fixed_wing" / "summary.txt")
    px4_summary = read_kv(RESULTS / "20260318_225551_px4_sih" / "summary.txt")

    success_preview = RESULTS / "full_docking_success_demo_20260315_preview.png"
    rviz_preview = RESULTS / "rviz_stable_after_fix_20260315_preview.png"
    mock_xy = RESULTS / "20260318_120407_mock_fixed_wing" / "trajectory_xy.png"
    mock_phase = RESULTS / "20260318_120407_mock_fixed_wing" / "phase_timeline.png"
    mock_speed = RESULTS / "20260318_120407_mock_fixed_wing" / "speed_profile.png"
    px4_wait = RESULTS / "20260318_225551_px4_sih" / "fixed_wing_wait_diagnostics.png"

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_NAVY
    add_textbox(
        slide, Inches(0.7), Inches(0.8), Inches(6.2), Inches(1.1),
        "Method\n面向异构空中航母任务的分阶段协同对接框架",
        size=24, bold=True, color=COLOR_WHITE
    )
    add_textbox(
        slide, Inches(0.75), Inches(2.1), Inches(5.7), Inches(0.9),
        "核心关键词：\n窗口触发  |  轨道切入  |  相对速度匹配  |  终端捕获",
        size=16, color=RGBColor(220, 228, 240)
    )
    add_metric_card(slide, Inches(0.75), Inches(3.35), Inches(1.8), Inches(1.0), "任务形态", "异构平台",
                    accent=COLOR_BLUE, note="carrier + fixed-wing")
    add_metric_card(slide, Inches(2.75), Inches(3.35), Inches(1.8), Inches(1.0), "目标速度级别", "10 m/s",
                    accent=COLOR_ORANGE, note="固定翼等待轨道")
    add_metric_card(slide, Inches(4.75), Inches(3.35), Inches(1.8), Inches(1.0), "终端目标", "0.2 m",
                    accent=COLOR_GREEN, note="dock window")
    add_textbox(
        slide, Inches(0.8), Inches(5.15), Inches(5.8), Inches(1.0),
        "不是看见目标就追，而是先判断可追窗口，\n再沿 fixed-wing 轨道切入，最后做末段速度匹配与捕获。",
        size=15, color=COLOR_WHITE
    )
    add_image(slide, success_preview, Inches(7.2), Inches(0.8), width=Inches(5.5))
    add_footer(slide, "Method deck | 2026-03-23 | 适合插入组会主报告的“方法”部分")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "1. 问题定义：不是双多旋翼靠近，而是异构空中回收", "carrier 为主动方，fixed-wing 在等待阶段保持高效盘旋，只在末段有限配合")
    add_bullets(
        slide, Inches(0.72), Inches(1.45), Inches(4.15), Inches(4.9),
        [
            "异构平台：carrier 是可悬停多旋翼，mini 是必须保持空速的 fixed-wing。",
            "任务形态：fixed-wing 先爬升到目标高度，再以定高、定半径、10 m/s 级空速等待盘旋。",
            "核心目标：carrier 主动去接机，末段让相对位置和相对速度同时收敛。",
            "因此本问题不只是“追上目标”，而是“在真实轨迹约束下稳稳捕获目标”。"
        ],
        title="场景假设", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.5
    )
    add_textbox(
        slide, Inches(5.15), Inches(1.55), Inches(3.0), Inches(1.2),
        "相对状态\nr = p_mini - p_carrier\nv_rel = v_mini - v_carrier",
        size=18, bold=True, fill_color=COLOR_LIGHT_BLUE, line_color=COLOR_BLUE,
        align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE
    )
    add_textbox(
        slide, Inches(5.15), Inches(2.95), Inches(3.0), Inches(1.2),
        "对接目标\nr -> r_d = [0, 0, 0.2]^T\nv_rel -> 0",
        size=18, bold=True, fill_color=RGBColor(227, 243, 232), line_color=COLOR_GREEN,
        align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE
    )
    add_textbox(
        slide, Inches(5.15), Inches(4.35), Inches(3.0), Inches(1.55),
        "关键区别\n不是最小化瞬时距离\n而是构造可捕获窗口",
        size=18, bold=True, fill_color=RGBColor(253, 239, 219), line_color=COLOR_ORANGE,
        align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE
    )
    add_image(slide, rviz_preview, Inches(8.45), Inches(1.45), width=Inches(4.25))
    add_footer(slide, "对外表述建议：这是“面向 10 m/s 级 fixed-wing 等待轨道”的异构空中回收问题。")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "2. 窗口触发 + 轨道切入：先决定什么时候追，再决定从哪里切入", "直接追飞机本体容易追尾和擦肩；我们先在等待轨道上找“可追窗口”")
    add_bullets(
        slide, Inches(0.72), Inches(1.45), Inches(4.5), Inches(2.25),
        [
            "Auto-start 不靠人工 START，也不只看距离阈值。",
            "先对等待轨道上的候选状态做几何筛选：distance, rel_x, rel_y, rel_vx, rel_vy, rel_vz。",
            "满足窗口约束后，再按评分选择最优候选窗口。"
        ],
        title="窗口触发思想", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.0
    )
    add_textbox(
        slide, Inches(0.82), Inches(3.95), Inches(4.35), Inches(1.85),
        "窗口评分示意\nJ_orbit = |d-d*| + 0.6|rel_x-x*| + 0.8|rel_y-y*|\n          + 0.3|rel_vy| + 0.8|rel_vz| + v_x 罚项",
        size=17, bold=True, fill_color=COLOR_LIGHT_BLUE, line_color=COLOR_BLUE
    )
    add_textbox(
        slide, Inches(0.82), Inches(5.95), Inches(4.35), Inches(0.55),
        "解释：核心不是“最近”，而是“这个窗口 carrier 能切进去，后面还接得住”。",
        size=13, color=COLOR_MUTED
    )
    add_textbox(
        slide, Inches(5.45), Inches(1.45), Inches(3.0), Inches(1.4),
        "轨道切入\np_int = p_mini + v_mini T_p - r_d",
        size=18, bold=True, fill_color=RGBColor(227, 243, 232), line_color=COLOR_GREEN,
        align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE
    )
    add_bullets(
        slide, Inches(5.45), Inches(3.05), Inches(3.0), Inches(2.75),
        [
            "不是朝 fixed-wing 的当前点硬追。",
            "而是先求一个 preview intercept point。",
            "carrier 先并入 fixed-wing 的轨道几何，再进入 close tracking。"
        ],
        title="轨道切入含义", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=13.8
    )
    add_image(slide, mock_xy, Inches(8.65), Inches(1.35), width=Inches(3.95))
    add_footer(slide, "术语一句话：先进入目标轨迹，再进入目标附近，而不是一开始就追目标本体。")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "3. 状态机匹配控制律：每个阶段做不同的事", "创新点不在单一控制器，而在于把“何时追/如何切入/如何收口”统一到了阶段式协同框架里")
    phase_specs = [
        ("SEARCH", COLOR_LIGHT_BLUE, [
            "目标：等待轨道观测 + 形成可追窗口",
            "动作：低增益位置保持，不盲追",
            "输出：最优窗口时刻与候选切入几何",
        ]),
        ("APPROACH", RGBColor(227, 243, 232), [
            "目标：切进 fixed-wing 轨道附近",
            "方法：LOS / intercept guidance",
            "命令：v_c^app = v_m + v_close + K_e e",
        ]),
        ("TRACKING", RGBColor(253, 239, 219), [
            "目标：同时压相对位置和相对速度",
            "方法：preview cut-in + relative guidance",
            "命令：v_c^trk = v_m + K_p e + K_v v_rel",
        ]),
        ("DOCKING", RGBColor(231, 238, 250), [
            "目标：从“看见窗口”变成“可捕获窗口”",
            "方法：track frame velocity shaping + backstepping",
            "命令：v_c^dock = v_m + Δv_match + αu_bs",
        ]),
    ]
    x = 0.68
    for title, color, bullets in phase_specs:
        add_bullets(slide, Inches(x), Inches(1.5), Inches(3.0), Inches(3.45), bullets,
                    title=title, fill_color=color, line_color=COLOR_BLUE, bullet_size=13.2)
        x += 3.12
    add_image(slide, mock_phase, Inches(0.95), Inches(5.15), width=Inches(6.15))
    add_bullets(
        slide, Inches(7.45), Inches(5.08), Inches(5.0), Inches(1.55),
        [
            "SEARCH 解决“什么时候开始追”",
            "APPROACH 解决“从哪里切进去”",
            "TRACKING 解决“怎么把几何和速度都压小”",
            "DOCKING 解决“最后几十厘米如何稳稳接住”",
        ],
        title="这就是方法的层次性", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=13.2
    )
    add_footer(slide, "这页可以直接回答老师：状态机不是工程拼接，而是控制目标与控制律的分层匹配。")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "4. 终端对接：末段优先匹配速度，再压位置误差", "真实空中对接的难点不是最近点，而是最近点附近能不能形成低相对速度捕获包络")
    add_bullets(
        slide, Inches(0.72), Inches(1.45), Inches(4.35), Inches(2.55),
        [
            "在 fixed-wing 航迹坐标系下分解误差：along-track / lateral / vertical。",
            "根据 terminal distance 自适应收紧相对速度上限，而不是一上来就强压 0 误差。",
            "用 backstepping 输入做辅助阻尼，但主目标是 velocity shaping。"
        ],
        title="终端控制核心", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=13.8
    )
    add_textbox(
        slide, Inches(0.82), Inches(4.25), Inches(4.2), Inches(1.7),
        "终端控制抽象\nΔv_match = t·f(e_a, v_a) + n·g(e_l, v_l) + k·h(e_z, v_z)\n捕获条件：|e| 小、|v_rel| 小，并保持若干控制周期",
        size=16.5, bold=True, fill_color=COLOR_LIGHT_BLUE, line_color=COLOR_BLUE
    )
    add_image(slide, mock_speed, Inches(5.45), Inches(1.45), width=Inches(3.45))
    add_image(slide, success_preview, Inches(9.1), Inches(1.45), width=Inches(3.05))
    add_metric_card(slide, Inches(5.55), Inches(5.35), Inches(1.8), Inches(0.95), "mock 最优", f"{mock_summary['min_distance_m']} m",
                    accent=COLOR_GREEN, note=mock_summary["final_phase"])
    add_metric_card(slide, Inches(7.5), Inches(5.35), Inches(1.8), Inches(0.95), "PX4 当前", f"{px4_summary['min_distance_m']} m",
                    accent=COLOR_ORANGE, note=px4_summary["best_window_phase"])
    add_metric_card(slide, Inches(9.45), Inches(5.35), Inches(2.0), Inches(0.95), "PX4 末段相对速度", f"{px4_summary['best_window_rel_speed_mps']} m/s",
                    accent=COLOR_RED, note="当前主瓶颈")
    add_footer(slide, "一句话：DOCKING 不是“只压距离”，而是“先把速度对齐，再把距离收进去”。")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "5. 为什么这套方法比直接追更有“研究味”", "它不是单点控制器，而是一种更接近空中航母任务形态的协同控制框架")
    add_textbox(
        slide, Inches(0.72), Inches(1.42), Inches(2.8), Inches(0.45), "常见直接追踪/平台着陆", size=16,
        bold=True, fill_color=RGBColor(245, 229, 226), line_color=COLOR_RED,
        align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE
    )
    add_textbox(
        slide, Inches(3.7), Inches(1.42), Inches(3.0), Inches(0.45), "本工作的方法表述", size=16,
        bold=True, fill_color=RGBColor(228, 238, 252), line_color=COLOR_BLUE,
        align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE
    )
    add_bullets(
        slide, Inches(0.72), Inches(1.95), Inches(2.8), Inches(3.65),
        [
            "看到目标就追",
            "常按距离阈值触发",
            "单一控制律全程追本体",
            "更适合同质平台或移动平台着陆",
        ],
        fill_color=COLOR_WHITE, line_color=RGBColor(230, 210, 205), bullet_size=14.0
    )
    add_bullets(
        slide, Inches(3.7), Inches(1.95), Inches(3.0), Inches(3.65),
        [
            "先做窗口评估，再触发追接",
            "先做轨道切入，再做 close tracking",
            "再用终端速度匹配完成捕获",
            "适配 carrier 主动 + fixed-wing 末端协同",
        ],
        fill_color=COLOR_WHITE, line_color=RGBColor(215, 226, 244), bullet_size=14.0
    )
    add_bullets(
        slide, Inches(7.0), Inches(1.42), Inches(5.45), Inches(4.15),
        [
            "独特性 1：把“何时开始追”显式建模成窗口触发问题。",
            "独特性 2：把“如何切入轨道”作为核心几何问题，而不是只做 point pursuit。",
            "独特性 3：把末段控制重心从“压距离”转到“速度匹配 + 捕获窗口”。",
            "独特性 4：面向异构平台，carrier 和 fixed-wing 采用非对称协同角色分工。",
            "工程外延：这套表述天然适合继续加 QP / MPC / 学习型窗口评估器。"
        ],
        title="你可以对老师强调的研究点", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=13.7
    )
    add_textbox(
        slide, Inches(0.95), Inches(6.05), Inches(11.5), Inches(0.55),
        "推荐收束句：本工作的创新不只是“把系统搭起来”，而是提出了一套面向高速 fixed-wing 目标的窗口驱动、轨道切入、分阶段协同对接框架。",
        size=15.5, bold=True, color=COLOR_NAVY, fill_color=RGBColor(232, 239, 251),
        line_color=COLOR_BLUE, align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE
    )
    add_footer(slide, "如果后面继续深挖，最自然的增强点就是把 DOCKING 末段升级成 QP/MPC，而前面三阶段保持现在的分层结构。")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "6. 从 method 到论文化表达：现在就能怎么讲，后续还能怎么加强", "这页适合答辩时用来承接老师的“方法还可以再硬核一点吗”")
    add_bullets(
        slide, Inches(0.72), Inches(1.5), Inches(5.7), Inches(4.9),
        [
            "当前已经能成立的表述：窗口驱动触发 + 轨道切入 + 分阶段协同对接。",
            "当前已经有实现支撑的要点：orbit-phase 评分、LOS/intercept、relative guidance、terminal velocity shaping。",
            "当前与真实 PX4 对接的难点：fixed-wing TECS / loiter / glide 约束使末段速度匹配变难。",
            "因此，方法价值已经成立；现在卡住的是“把真实飞控链上的末段窗口压稳”。"
        ],
        title="现在就能讲得住的内容", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.0
    )
    add_bullets(
        slide, Inches(6.75), Inches(1.5), Inches(5.55), Inches(4.9),
        [
            "硬核升级方向 1：把 orbit window score 学习化或优化化。",
            "硬核升级方向 2：在 DOCKING 阶段加入 QP / MPC，仅负责终端约束满足。",
            "硬核升级方向 3：把 fixed-wing TECS 状态引入 terminal controller，形成 TECS-aware docking。",
            "硬核升级方向 4：把相对状态与视线角融合成统一 guidance law。"
        ],
        title="下一版更像论文方法的增强方向", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.0
    )
    add_image(slide, px4_wait, Inches(4.25), Inches(5.35), width=Inches(4.8))
    add_footer(slide, "建议讲法：现在的方法骨架已经足够清楚，后续增强是在这个骨架上增加更强的终端优化与飞控感知。")

    return prs


def main() -> None:
    DOCS.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(OUTPUT_PPTX)
    prs.save(OUTPUT_PPTX_CN)
    print(OUTPUT_PPTX)
    print(OUTPUT_PPTX_CN)


if __name__ == "__main__":
    main()
