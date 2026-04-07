#!/usr/bin/env python3

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
RESULTS = ROOT / "results"
DOCS = ROOT / "docs"
OUTPUT_PPTX = DOCS / "group_meeting_progress_20260323.pptx"
OUTPUT_PPTX_CN = DOCS / "组会汇报_空中对接进展_20260323.pptx"

COLOR_BG = RGBColor(246, 248, 252)
COLOR_NAVY = RGBColor(15, 30, 60)
COLOR_BLUE = RGBColor(35, 94, 184)
COLOR_LIGHT_BLUE = RGBColor(218, 232, 252)
COLOR_GREEN = RGBColor(29, 125, 72)
COLOR_ORANGE = RGBColor(205, 121, 35)
COLOR_RED = RGBColor(180, 55, 42)
COLOR_TEXT = RGBColor(32, 37, 43)
COLOR_MUTED = RGBColor(96, 106, 120)
COLOR_WHITE = RGBColor(255, 255, 255)


def read_kv(path: Path) -> dict[str, str]:
    data: dict[str, str] = {}
    for line in path.read_text(encoding="utf-8", errors="ignore").splitlines():
        if "=" in line:
            key, value = line.split("=", 1)
            data[key.strip()] = value.strip()
    return data


def fmt_run_metrics(run_dir: Path) -> dict[str, str]:
    summary = read_kv(run_dir / "summary.txt")
    metadata = read_kv(run_dir / "metadata.txt") if (run_dir / "metadata.txt").exists() else {}
    out = {**summary, **metadata}
    return out


def set_text_style(run, size: float, bold: bool = False, color: RGBColor = COLOR_TEXT, font_name: str = "Microsoft YaHei") -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_textbox(slide, left, top, width, height, text="", size=20, bold=False, color=COLOR_TEXT,
                fill_color: RGBColor | None = None, line_color: RGBColor | None = None,
                align=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.TOP, margin=0.10):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = vertical
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_text_style(run, size=size, bold=bold, color=color)

    if fill_color is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()

    if line_color is not None:
        box.line.color.rgb = line_color
    else:
        box.line.fill.background()
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(10.8), Inches(0.6), title, size=26, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.62), Inches(0.92), Inches(9.8), Inches(0.35), subtitle, size=10.5, color=COLOR_MUTED)


def add_bullets(slide, left, top, width, height, bullets: Iterable[str], title: str | None = None,
                fill_color: RGBColor | None = None, line_color: RGBColor | None = None,
                title_color: RGBColor = COLOR_BLUE, bullet_size: float = 18.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)

    if fill_color is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()
    if line_color is not None:
        box.line.color.rgb = line_color
    else:
        box.line.fill.background()

    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        set_text_style(run, 16, bold=True, color=title_color)
    else:
        tf.clear()

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            set_text_style(p.runs[0], bullet_size, color=COLOR_TEXT)
    return box


def add_image(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_metric_card(slide, left, top, width, height, title: str, value: str,
                    accent: RGBColor = COLOR_BLUE, note: str | None = None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.10)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_text_style(r1, 12, bold=True, color=accent)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    set_text_style(r2, 24, bold=True, color=COLOR_TEXT)
    if note:
        p3 = tf.add_paragraph()
        r3 = p3.add_run()
        r3.text = note
        set_text_style(r3, 10, color=COLOR_MUTED)
    return shape


def add_footer(slide, text: str) -> None:
    add_textbox(slide, Inches(0.6), Inches(6.95), Inches(11.8), Inches(0.2), text, size=9, color=COLOR_MUTED)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    mock_best = fmt_run_metrics(RESULTS / "20260318_120407_mock_fixed_wing")
    px4_poc = fmt_run_metrics(RESULTS / "20260317_222631_px4_sih")
    px4_best = fmt_run_metrics(RESULTS / "20260318_225551_px4_sih")
    px4_speed_variant = fmt_run_metrics(RESULTS / "20260318_234151_px4_sih")

    cover_img = RESULTS / "full_docking_success_demo_20260315_preview.png"
    rviz_img = RESULTS / "rviz_stable_after_fix_20260315_preview.png"
    mock_xy = RESULTS / "20260318_120407_mock_fixed_wing" / "trajectory_xy.png"
    mock_dist = RESULTS / "20260318_120407_mock_fixed_wing" / "distance_convergence.png"
    poc_xy = RESULTS / "20260317_222631_px4_sih" / "trajectory_xy.png"
    poc_dist = RESULTS / "20260317_222631_px4_sih" / "distance_convergence.png"
    best_xy = RESULTS / "20260318_225551_px4_sih" / "trajectory_xy.png"
    best_dist = RESULTS / "20260318_225551_px4_sih" / "distance_convergence.png"
    best_wait = RESULTS / "20260318_225551_px4_sih" / "fixed_wing_wait_diagnostics.png"

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_NAVY
    add_textbox(slide, Inches(0.7), Inches(0.8), Inches(6.1), Inches(1.2),
                "空中对接系统复现与 PX4 迁移进展", size=28, bold=True, color=COLOR_WHITE)
    add_textbox(slide, Inches(0.72), Inches(1.95), Inches(5.8), Inches(0.6),
                "基于论文方法的阶段式对接控制\nROS 2 通信 + RViz 可视化 + PX4 SITL 迁移", size=16, color=RGBColor(220, 228, 240))
    add_metric_card(slide, Inches(0.75), Inches(3.15), Inches(1.8), Inches(1.1), "mock fixed-wing", "0.2 m",
                    accent=COLOR_GREEN, note="COMPLETED")
    add_metric_card(slide, Inches(2.75), Inches(3.15), Inches(1.8), Inches(1.1), "PX4 SIH 早期 POC", "0.8 m",
                    accent=COLOR_ORANGE, note="COMPLETED")
    add_metric_card(slide, Inches(4.75), Inches(3.15), Inches(1.8), Inches(1.1), "当前真实链最好", "4.237 m",
                    accent=COLOR_BLUE, note="DOCKING")
    add_textbox(slide, Inches(0.78), Inches(5.15), Inches(5.8), Inches(1.1),
                "一句话总结：简单 docking 已打通，当前主要工作是把论文式分阶段追踪迁移到真实 PX4 fixed-wing 链路，并把末段从 4 m 级继续压到 0.2 m。",
                size=15, color=COLOR_WHITE)
    add_image(slide, cover_img, Inches(7.2), Inches(0.75), width=Inches(5.6))
    add_footer(slide, "组会汇报 | 2026-03-23 | easydocking")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "项目路线与当前位置", "从论文方法理解，到简单 docking 打通，再到 PX4 真实飞控迁移")
    stages = [
        ("S1", "简化原型", "双四轴/简化模型\n打通状态机、RViz、日志"),
        ("S2", "mock fixed-wing", "carrier 主动追接\nfixed-wing 末段配合减速"),
        ("S3", "PX4 SIH 早期 POC", "桥接/Offboard/日志闭环\n0.8 m COMPLETED"),
        ("S4", "当前真实链", "30 m 定高盘旋 + 自动触发\nbest 4.237 m DOCKING"),
    ]
    x = 0.7
    colors = [COLOR_LIGHT_BLUE, RGBColor(227, 243, 232), RGBColor(253, 239, 219), RGBColor(231, 238, 250)]
    for idx, (tag, name, desc) in enumerate(stages):
        add_textbox(slide, Inches(x), Inches(1.85), Inches(2.9), Inches(1.6),
                    f"{tag}\n{name}\n{desc}", size=16, bold=False, fill_color=colors[idx],
                    line_color=COLOR_BLUE if idx == 3 else COLOR_MUTED, align=PP_ALIGN.CENTER,
                    vertical=MSO_ANCHOR.MIDDLE)
        if idx < len(stages) - 1:
            add_textbox(slide, Inches(x + 2.95), Inches(2.4), Inches(0.35), Inches(0.4), "→",
                        size=24, bold=True, color=COLOR_BLUE, align=PP_ALIGN.CENTER,
                        vertical=MSO_ANCHOR.MIDDLE)
        x += 3.15
    add_bullets(
        slide, Inches(0.75), Inches(4.1), Inches(6.0), Inches(2.2),
        [
            "总目标：把论文中的空中对接思想复现为可运行、可观察、可复现实验系统。",
            "平台目标：PX4 负责飞控，ROS 2 负责通信与控制编排，RViz 负责在线可视化。",
            "当前位置：简单 docking 已完成，核心难点转为真实 fixed-wing PX4 末段收口。"
        ],
        title="本阶段已经完成什么", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235)
    )
    add_image(slide, rviz_img, Inches(7.15), Inches(4.05), width=Inches(5.35))
    add_footer(slide, "当前主线：在保持 30 m / 10 m·s⁻¹ 等待轨道约束下，继续压真实 PX4 末段误差。")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "系统架构与代码落点", "不仅是跑通仿真，更是把控制、飞控、可视化、报告链路接成一个闭环")
    box_specs = [
        (Inches(0.7), Inches(1.65), Inches(2.3), Inches(1.05), "Carrier\nPX4 SITL / Offboard"),
        (Inches(3.5), Inches(1.65), Inches(2.6), Inches(1.05), "ROS 2 协同控制层\nDockingController"),
        (Inches(6.65), Inches(1.65), Inches(2.45), Inches(1.05), "Mini\nmock fixed-wing / PX4 SIH"),
        (Inches(9.6), Inches(1.65), Inches(2.7), Inches(1.05), "RViz + Logger + Report\n可视化 / CSV / 图表"),
    ]
    for left, top, width, height, text in box_specs:
        add_textbox(slide, left, top, width, height, text, size=16, bold=True,
                    fill_color=COLOR_WHITE, line_color=COLOR_BLUE, align=PP_ALIGN.CENTER,
                    vertical=MSO_ANCHOR.MIDDLE)
    for arrow_x in [3.05, 6.18, 9.18]:
        add_textbox(slide, Inches(arrow_x), Inches(1.98), Inches(0.28), Inches(0.3), "→",
                    size=22, bold=True, color=COLOR_BLUE, align=PP_ALIGN.CENTER,
                    vertical=MSO_ANCHOR.MIDDLE)
    add_bullets(
        slide, Inches(0.75), Inches(3.15), Inches(5.9), Inches(2.7),
        [
            "核心控制器：[docking_controller.cpp]：状态机、相对位姿、末段速度整形。",
            "fixed-wing 桥接：[px4_fixed_wing_bridge.py]：等待轨道、末段 glide、速度调度。",
            "自动触发：[auto_start_docking.py]：orbit-phase score-based auto-start。",
            "实验闭环：[experiment_logger.py] + [run_px4_sih_docking_experiment.sh]：一键跑、自动出图。"
        ],
        title="关键模块", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235)
    )
    add_bullets(
        slide, Inches(7.0), Inches(3.15), Inches(5.35), Inches(2.7),
        [
            "控制链：状态机输出 → bridge/odometry → PX4 Offboard / orbit / TECS。",
            "观测链：ROS topic / docking status / relative pose / trajectory markers。",
            "复现实验链：同一套 runner 输出 trajectory、distance、phase、speed、wait diagnostics。"
        ],
        title="为什么这一步重要", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235)
    )
    add_footer(slide, "重点不是单点算法，而是把“算法 + 飞控 + 可视化 + 日志”接成了可复现实验平台。")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "核心追踪算法：分阶段而不是一把梭", "每个阶段目标不同，所以控制律、角色分工、触发条件都不同")
    phase_blocks = [
        ("SEARCH", COLOR_LIGHT_BLUE, [
            "fixed-wing 先上升到 30 m 并定高盘旋",
            "不盲追目标，而是先等待“可追窗口”形成",
            "auto-start 依据 orbit-phase 几何评分自动触发",
        ]),
        ("APPROACH", RGBColor(227, 243, 232), [
            "carrier 用 LOS / intercept guidance 拉近几何",
            "重点收横向偏差和相对高度",
            "fixed-wing 仍保持等待轨道为主",
        ]),
        ("TRACKING", RGBColor(253, 239, 219), [
            "切到相对位置/速度联合制导",
            "让 carrier 主动切到 fixed-wing 航迹前方",
            "同时为 glide / terminal cooperation 创造条件",
        ]),
        ("DOCKING", RGBColor(231, 238, 250), [
            "末段目标从“看得见”变成“可捕获”",
            "carrier 做 terminal velocity shaping + backstepping",
            "fixed-wing 进入 glide 并末段减速配合",
        ]),
    ]
    x = 0.65
    for title, color, bullets in phase_blocks:
        add_bullets(slide, Inches(x), Inches(1.55), Inches(3.0), Inches(4.8), bullets,
                    title=title, fill_color=color, line_color=COLOR_BLUE, bullet_size=14.5)
        x += 3.15
    add_footer(slide, "核心思想：先制造窗口，再做终端捕获。这样比全程单一 PID/直接追最近点更稳定，也更容易迁移到 PX4。")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "独特性与优势：不是只做工程堆叠", "针对老师的问题，这页直接回答“你的工作与常见方法有什么本质不同”")
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(2.7), Inches(0.45), "常见直接追踪方案", size=16, bold=True,
                fill_color=RGBColor(245, 229, 226), line_color=COLOR_RED, align=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(3.55), Inches(1.5), Inches(2.8), Inches(0.45), "我的方案", size=16, bold=True,
                fill_color=RGBColor(228, 238, 252), line_color=COLOR_BLUE, align=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE)
    compare_left = [
        "人工发 START，或者只按距离阈值触发",
        "单一控制律全程追飞机本体",
        "默认双方都一直积极机动",
        "mock 仿真可跑，但很难迁到真实飞控",
    ]
    compare_right = [
        "orbit-phase score-based auto-start，先找可追窗口",
        "SEARCH / APPROACH / TRACKING / DOCKING 分层控制",
        "carrier 主动追接，fixed-wing 只在末段减速配合",
        "从 mock 直接迁到 PX4 SIH，并保留日志/报告闭环",
    ]
    add_bullets(slide, Inches(0.7), Inches(2.0), Inches(2.7), Inches(3.3), compare_left,
                fill_color=COLOR_WHITE, line_color=RGBColor(230, 210, 205), bullet_size=14.5)
    add_bullets(slide, Inches(3.55), Inches(2.0), Inches(2.8), Inches(3.3), compare_right,
                fill_color=COLOR_WHITE, line_color=RGBColor(215, 226, 244), bullet_size=14.5)
    add_bullets(
        slide, Inches(6.7), Inches(1.5), Inches(5.8), Inches(4.0),
        [
            "独特点 1：把“何时开始追”变成了核心问题，而不是只调追踪器增益。",
            "独特点 2：把异构平台约束显式写进状态机和速度调度里。",
            "独特点 3：把论文式分阶段思想落成了可复现的 PX4/ROS2 工程框架。",
            "优势：更可解释、更易调参、更容易从仿真迁移到真实飞控。"
        ],
        title="可以在组会上强调的点", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235)
    )
    add_footer(slide, "一句话：我的工作不是“搭个平台”，而是把论文里的分阶段思想做成了面向异构飞行器的可迁移方法框架。")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "可提炼成研究贡献的 3 点", "如果老师追问“独特之处是什么”，建议优先从下面三点来讲")
    add_bullets(
        slide, Inches(0.7), Inches(1.45), Inches(3.9), Inches(4.65),
        [
            "不是看见目标就追，而是先判断等待轨道上是否形成“可追窗口”。",
            "触发逻辑不再靠人工 START 或单一距离阈值，而是看 orbit phase、相对几何和可达性评分。",
            "好处：减少无效追逐，提升进入 TRACKING / DOCKING 的概率。"
        ],
        title="贡献 1：窗口驱动的追踪触发", fill_color=COLOR_WHITE, line_color=COLOR_BLUE, bullet_size=14.5
    )
    add_bullets(
        slide, Inches(4.75), Inches(1.45), Inches(3.9), Inches(4.65),
        [
            "SEARCH、APPROACH、TRACKING、DOCKING 各自目标不同，对应不同控制律。",
            "SEARCH 强调窗口形成，APPROACH 用 LOS / intercept 收几何，TRACKING 压位置和速度，DOCKING 再做终端速度匹配。",
            "好处：比单一 PID 或一路 pure pursuit 更符合任务结构，也更容易调参。"
        ],
        title="贡献 2：状态机匹配控制律", fill_color=COLOR_WHITE, line_color=COLOR_BLUE, bullet_size=14.5
    )
    add_bullets(
        slide, Inches(8.8), Inches(1.45), Inches(3.8), Inches(4.65),
        [
            "carrier 是主动方，负责大部分机动；fixed-wing 在等待阶段尽量保持稳定盘旋。",
            "只有在 APPROACH / DOCKING 末段才让 fixed-wing 通过降速 / glide 配合收口。",
            "好处：尊重异构平台物理约束，避免双方同时剧烈机动导致窗口失稳。"
        ],
        title="贡献 3：异构平台角色分工", fill_color=COLOR_WHITE, line_color=COLOR_BLUE, bullet_size=14.5
    )
    add_footer(slide, "这三点合起来，才是这项工作的“方法特色”：窗口先行、分层制导、异构协同。")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "阶段成果 I：mock fixed-wing 协同降落已经打通", "证明算法和状态机本身是通的，不是只能停留在接近而不能闭环")
    add_metric_card(slide, Inches(0.75), Inches(1.45), Inches(1.75), Inches(1.0), "最小距离", f"{mock_best['min_distance_m']} m",
                    accent=COLOR_GREEN, note="run: 20260318_120407")
    add_metric_card(slide, Inches(2.65), Inches(1.45), Inches(1.75), Inches(1.0), "最终阶段", mock_best["final_phase"],
                    accent=COLOR_GREEN, note="已完成对接")
    add_metric_card(slide, Inches(4.55), Inches(1.45), Inches(1.75), Inches(1.0), "末段相对速度", f"{mock_best['best_window_rel_speed_mps']} m/s",
                    accent=COLOR_GREEN, note="best window")
    add_bullets(
        slide, Inches(0.75), Inches(2.65), Inches(5.55), Inches(1.8),
        [
            "carrier 主动去追接 fixed-wing，fixed-wing 末段配合减速，最终进入 COMPLETED。",
            "说明：阶段状态机、终端捕获逻辑、rigid attach / completed 判据是成立的。",
            "这是向 PX4 迁移前的关键算法验证版本。"
        ],
        title="这个结果说明了什么", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.5
    )
    add_image(slide, mock_xy, Inches(6.55), Inches(1.35), width=Inches(3.0))
    add_image(slide, mock_dist, Inches(9.75), Inches(1.35), width=Inches(3.0))
    add_footer(slide, "代表结果：results/20260318_120407_mock_fixed_wing")

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "阶段成果 II：真实 PX4 SIH 早期版本曾打通 COMPLETED", "说明 ROS2-PX4 bridge、offboard、日志链、报告链和真实飞控接口已经接通过")
    add_metric_card(slide, Inches(0.75), Inches(1.45), Inches(1.75), Inches(1.0), "最小距离", f"{px4_poc['min_distance_m']} m",
                    accent=COLOR_ORANGE, note="run: 20260317_222631")
    add_metric_card(slide, Inches(2.65), Inches(1.45), Inches(1.75), Inches(1.0), "最终阶段", px4_poc["final_phase"],
                    accent=COLOR_ORANGE, note="真实 PX4 SIH")
    add_metric_card(slide, Inches(4.55), Inches(1.45), Inches(1.75), Inches(1.0), "best window", px4_poc["best_window_phase"],
                    accent=COLOR_ORANGE, note="rigid attach POC")
    add_bullets(
        slide, Inches(0.75), Inches(2.65), Inches(5.55), Inches(1.8),
        [
            "这不是最终版本，但说明真实 PX4 flight stack 不是“完全接不通”。",
            "完成了 PX4 SITL、MicroXRCEAgent、ROS 2 bridge、RViz、自动报告的首次闭环。",
            "后面为了更真实的 fixed-wing 定高盘旋与速度约束，任务难度比这个版本更高。"
        ],
        title="为什么还要继续做", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.5
    )
    add_image(slide, poc_xy, Inches(6.55), Inches(1.35), width=Inches(3.0))
    add_image(slide, poc_dist, Inches(9.75), Inches(1.35), width=Inches(3.0))
    add_footer(slide, "代表结果：results/20260317_222631_px4_sih")

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "当前状态：更真实 fixed-wing 约束下已稳定进入 DOCKING", "现在正在做的，不是“能不能跑”，而是“能不能在真实 PX4 约束下稳定把末段收进去”")
    add_metric_card(slide, Inches(0.75), Inches(1.35), Inches(1.75), Inches(1.0), "最小距离", f"{px4_best['min_distance_m']} m",
                    accent=COLOR_BLUE, note="run: 20260318_225551")
    add_metric_card(slide, Inches(2.65), Inches(1.35), Inches(1.75), Inches(1.0), "best window", px4_best["best_window_phase"],
                    accent=COLOR_BLUE, note=f"rel speed {px4_best['best_window_rel_speed_mps']} m/s")
    add_metric_card(slide, Inches(4.55), Inches(1.35), Inches(1.75), Inches(1.0), "等待轨道", "30 m / 80 m / 10 m·s⁻¹",
                    accent=COLOR_BLUE, note="当前真实链约束")
    add_bullets(
        slide, Inches(0.75), Inches(2.55), Inches(5.55), Inches(2.4),
        [
            "fixed-wing 可在 30 m 左右进入等待轨道，auto-start 和 offboard glide 都已接通。",
            "best run：进入 DOCKING，最近点 4.237 m，best window 相对速度 5.699 m/s。",
            "近似速度更优的候选组：20260318_234151，4.466 m / 5.625 m·s⁻¹。",
            "结论：简单 docking 已经完成，当前真正卡的是真实 fixed-wing 末段速度匹配。"
        ],
        title="现在已经稳定下来的能力", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.0
    )
    add_image(slide, best_xy, Inches(6.45), Inches(1.3), width=Inches(2.2))
    add_image(slide, best_dist, Inches(8.85), Inches(1.3), width=Inches(3.7))
    add_image(slide, best_wait, Inches(6.45), Inches(4.05), width=Inches(6.1))
    add_footer(slide, "代表结果：results/20260318_225551_px4_sih | 近似速度更优候选：results/20260318_234151_px4_sih")

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "当前瓶颈与下一步", "简单 docking 已经结束；下一阶段是把真实 PX4 fixed-wing 版本从 4 m 级进一步压到 0.2 m 级")
    add_bullets(
        slide, Inches(0.75), Inches(1.55), Inches(5.7), Inches(4.6),
        [
            "当前主要瓶颈：DOCKING 阶段 along-track 速度仍偏大，属于“靠近了但没被稳稳接住”。",
            "fixed-wing 的 TECS / loiter / glide 约束使真实链比 mock 版本难得多。",
            "控制器大改容易退化，说明问题已从“能不能接通”转成“终端窗口怎么稳稳抓住”。",
        ],
        title="为什么 0.2 m 还没在真实 PX4 上实现", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235)
    )
    add_bullets(
        slide, Inches(6.75), Inches(1.55), Inches(5.75), Inches(4.6),
        [
            "继续围绕末段 along-track speed matching 做终端控制，而不是再大范围盲调。",
            "保持 fixed-wing 30 m 定高盘旋 + 10 m/s 约束，只在 terminal phase 再减速配合。",
            "沿当前最优链条继续：orbit-phase auto-start + staged glide + terminal capture。",
            "目标：真实 PX4 fixed-wing 链下，实现 <0.2 m 且相对速度基本一致的收口。"
        ],
        title="下一步准备怎么做", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235)
    )
    add_textbox(slide, Inches(0.85), Inches(6.2), Inches(11.6), Inches(0.55),
                "结论：这项工作的价值不只是把系统搭起来，而是构建了一套面向异构飞行器、可迁移到真实飞控的分阶段空中对接方法框架。",
                size=16, bold=True, color=COLOR_NAVY, fill_color=RGBColor(232, 239, 251),
                line_color=COLOR_BLUE, align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE)
    add_footer(slide, "汇报建议收束句：论文思想已经被我们做成了系统；下一步是把真实 PX4 末段从“能进 DOCKING”推进到“真稳定捕获”。")

    return prs


def main() -> None:
    DOCS.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(OUTPUT_PPTX)
    prs.save(OUTPUT_PPTX_CN)
    print(OUTPUT_PPTX)
    print(OUTPUT_PPTX_CN)


if __name__ == "__main__":
    main()
