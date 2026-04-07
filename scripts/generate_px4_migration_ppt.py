#!/usr/bin/env python3

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
RESULTS = ROOT / "results"
OUTPUT_PPTX = DOCS / "px4_migration_bottleneck_20260323.pptx"
OUTPUT_PPTX_CN = DOCS / "PX4迁移瓶颈分析_20260323.pptx"

COLOR_BG = RGBColor(246, 248, 252)
COLOR_NAVY = RGBColor(15, 30, 60)
COLOR_BLUE = RGBColor(35, 94, 184)
COLOR_LIGHT_BLUE = RGBColor(218, 232, 252)
COLOR_GREEN = RGBColor(29, 125, 72)
COLOR_ORANGE = RGBColor(205, 121, 35)
COLOR_RED = RGBColor(180, 55, 42)
COLOR_TEXT = RGBColor(32, 37, 43)
COLOR_MUTED = RGBColor(96, 106, 120)
COLOR_WHITE = RGBColor(255, 255, 255)


def read_kv(path: Path) -> dict[str, str]:
    data: dict[str, str] = {}
    for line in path.read_text(encoding="utf-8", errors="ignore").splitlines():
        if "=" in line:
            key, value = line.split("=", 1)
            data[key.strip()] = value.strip()
    return data


def set_text_style(run, size: float, bold: bool = False, color: RGBColor = COLOR_TEXT,
                   font_name: str = "Microsoft YaHei") -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_textbox(slide, left, top, width, height, text="", size=20, bold=False, color=COLOR_TEXT,
                fill_color: RGBColor | None = None, line_color: RGBColor | None = None,
                align=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.TOP, margin=0.10):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = vertical
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_text_style(run, size=size, bold=bold, color=color)

    if fill_color is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()

    if line_color is not None:
        box.line.color.rgb = line_color
    else:
        box.line.fill.background()
    return box


def add_bullets(slide, left, top, width, height, bullets: Iterable[str], title: str | None = None,
                fill_color: RGBColor | None = None, line_color: RGBColor | None = None,
                title_color: RGBColor = COLOR_BLUE, bullet_size: float = 18.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)

    if fill_color is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    else:
        box.fill.background()
    if line_color is not None:
        box.line.color.rgb = line_color
    else:
        box.line.fill.background()

    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        set_text_style(run, 16, bold=True, color=title_color)
    else:
        tf.clear()

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            set_text_style(p.runs[0], bullet_size, color=COLOR_TEXT)
    return box


def add_metric_card(slide, left, top, width, height, title: str, value: str,
                    accent: RGBColor = COLOR_BLUE, note: str | None = None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.10)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_text_style(r1, 12, bold=True, color=accent)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    set_text_style(r2, 22, bold=True, color=COLOR_TEXT)
    if note:
        p3 = tf.add_paragraph()
        r3 = p3.add_run()
        r3.text = note
        set_text_style(r3, 10, color=COLOR_MUTED)
    return shape


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(11.6), Inches(0.6), title, size=26, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.62), Inches(0.92), Inches(11.2), Inches(0.35), subtitle, size=10.5, color=COLOR_MUTED)


def add_footer(slide, text: str) -> None:
    add_textbox(slide, Inches(0.6), Inches(6.95), Inches(11.8), Inches(0.2), text, size=9, color=COLOR_MUTED)


def add_image(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    poc = read_kv(RESULTS / "20260317_222631_px4_sih" / "summary.txt")
    best = read_kv(RESULTS / "20260318_225551_px4_sih" / "summary.txt")
    metadata = read_kv(RESULTS / "20260318_225551_px4_sih" / "metadata.txt")

    rviz_preview = RESULTS / "rviz_stable_after_fix_20260315_preview.png"
    wait_diag = RESULTS / "20260318_225551_px4_sih" / "fixed_wing_wait_diagnostics.png"
    best_xy = RESULTS / "20260318_225551_px4_sih" / "trajectory_xy.png"
    best_dist = RESULTS / "20260318_225551_px4_sih" / "distance_convergence.png"

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_NAVY
    add_textbox(
        slide, Inches(0.7), Inches(0.82), Inches(6.3), Inches(1.0),
        "为什么迁移到 PX4 还没有完全打通？",
        size=27, bold=True, color=COLOR_WHITE
    )
    add_textbox(
        slide, Inches(0.74), Inches(1.95), Inches(5.95), Inches(1.15),
        "结论先行：\n不是 PX4 没接进来，而是接进来以后，fixed-wing 不再是理想目标，末段可捕获窗口被真实飞控约束显著压缩。",
        size=16, color=RGBColor(220, 228, 240)
    )
    add_metric_card(slide, Inches(0.78), Inches(3.45), Inches(1.9), Inches(1.0), "早期 PX4 POC", "0.8 m",
                    accent=COLOR_GREEN, note="COMPLETED")
    add_metric_card(slide, Inches(2.9), Inches(3.45), Inches(1.9), Inches(1.0), "当前真实链最好", "4.237 m",
                    accent=COLOR_ORANGE, note="best window = DOCKING")
    add_metric_card(slide, Inches(5.02), Inches(3.45), Inches(1.9), Inches(1.0), "末段相对速度", "5.699 m/s",
                    accent=COLOR_RED, note="当前主瓶颈")
    add_textbox(
        slide, Inches(0.82), Inches(5.15), Inches(6.05), Inches(0.95),
        "一句话：通信、飞控接口、状态机都通了，\n真正卡住的是“真实 fixed-wing 在 PX4 约束下还能不能被稳稳接住”。",
        size=15, color=COLOR_WHITE
    )
    add_image(slide, best_xy, Inches(7.15), Inches(0.9), width=Inches(5.4))
    add_footer(slide, "可直接作为组会里“问题定位”页的开场。")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "1. 已经做到什么，没做到什么", "先把“没接进去”这件事说清楚：现在的问题不是接口不通，而是末段捕获不稳")
    add_bullets(
        slide, Inches(0.72), Inches(1.45), Inches(5.5), Inches(2.65),
        [
            "已做到：ROS 2 <-> PX4 通信、MicroXRCE、SITL、RViz、自动日志和报告都已跑通。",
            "已做到：早期 PX4 SIH 版本曾达到 0.8 m / COMPLETED，证明 flight stack 接口不是完全接不通。",
            "已做到：当前更真实 fixed-wing 等待轨道链已能稳定进入 DOCKING。"
        ],
        title="已经实现的部分", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.3
    )
    add_bullets(
        slide, Inches(0.72), Inches(4.35), Inches(5.5), Inches(1.9),
        [
            "尚未做到：在 30 m 定高、10 m/s 级等待轨道约束下，稳定实现 0.2 m、低相对速度、一致航向的真实终端捕获。",
            "当前现象：能靠近、能进 DOCKING，但容易高速擦过，不能稳定守住 COMPLETED。"
        ],
        title="当前还没打通的部分", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.3
    )
    add_metric_card(slide, Inches(6.6), Inches(1.55), Inches(1.8), Inches(0.95), "POC", poc["final_phase"],
                    accent=COLOR_GREEN, note=f"min {poc['min_distance_m']} m")
    add_metric_card(slide, Inches(8.55), Inches(1.55), Inches(1.8), Inches(0.95), "当前最佳", best["best_window_phase"],
                    accent=COLOR_ORANGE, note=f"min {best['min_distance_m']} m")
    add_metric_card(slide, Inches(10.5), Inches(1.55), Inches(1.8), Inches(0.95), "run 结束", best["final_phase"],
                    accent=COLOR_RED, note=f"final {best['final_distance_m']} m")
    add_image(slide, best_dist, Inches(6.55), Inches(2.8), width=Inches(5.75))
    add_footer(slide, "推荐讲法：问题已经从“能不能跑”升级成“真实约束下能不能稳定收口”。")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "2. 第一层根因：接入 PX4 后，我们失去了对 fixed-wing 的直接控制权", "mock 中我们控制的是目标机状态；PX4 中我们控制的是飞控输入，最终状态由 PX4 自己决定")
    add_bullets(
        slide, Inches(0.72), Inches(1.45), Inches(4.0), Inches(4.8),
        [
            "mock 版本：我们几乎直接定义 mini 的轨迹、速度和末段减速。",
            "PX4 版本：bridge 只能给 DO_ORBIT、DO_CHANGE_SPEED、offboard setpoint 这类高层命令。",
            "因此 fixed-wing 的真实轨迹，要经过 PX4 内部导航、控制、TECS 再实现。",
            "结果就是：我们不是在“直接控飞机”，而是在“给飞控提建议”。"
        ],
        title="接口层变化", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.2
    )
    add_textbox(
        slide, Inches(5.05), Inches(1.6), Inches(3.0), Inches(1.15),
        "mock\nstate-level target motion",
        size=18, bold=True, fill_color=RGBColor(227, 243, 232), line_color=COLOR_GREEN,
        align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE
    )
    add_textbox(
        slide, Inches(5.05), Inches(3.0), Inches(3.0), Inches(1.15),
        "PX4\ncommand-level guidance",
        size=18, bold=True, fill_color=COLOR_LIGHT_BLUE, line_color=COLOR_BLUE,
        align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE
    )
    add_textbox(
        slide, Inches(5.05), Inches(4.42), Inches(3.0), Inches(1.45),
        "后果\n目标机行为不再完全按算法理想模型执行",
        size=18, bold=True, fill_color=RGBColor(253, 239, 219), line_color=COLOR_ORANGE,
        align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE
    )
    add_image(slide, rviz_preview, Inches(8.35), Inches(1.55), width=Inches(4.05))
    add_footer(slide, "一句话：迁移后不是控制器失效，而是控制对象从“理想目标”变成了“PX4 管理下的真实 fixed-wing”。")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "3. 第二层根因：fixed-wing 的真实物理与 TECS/loiter 约束会吃掉窗口", "固定翼必须服从空速、转弯半径、升力与能量管理，无法像多旋翼一样按需慢速等待")
    add_bullets(
        slide, Inches(0.72), Inches(1.45), Inches(5.1), Inches(2.7),
        [
            "当前真实链要求：30 m 左右定高、80 m 半径、10 m/s 级等待盘旋。",
            "但 fixed-wing 不能像多旋翼那样原地等，它必须维持空速并通过倾侧转弯。",
            "降速、转弯、定高三者天然耦合，所以等待轨道本身就比 mock 目标更“软”。"
        ],
        title="飞行器层约束", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.0
    )
    add_bullets(
        slide, Inches(0.72), Inches(4.45), Inches(5.1), Inches(1.75),
        [
            f"日志显示：alt_min={best['mini_altitude_min_m']} m, alt_max={best['mini_altitude_max_m']} m",
            f"tas_min={best['mini_tecs_min_tas_mps']} m/s, tas_max={best['mini_tecs_max_tas_mps']} m/s, underspeed={best['mini_tecs_max_underspeed_ratio']}",
            "这说明 fixed-wing 并没有长期稳定守住理想的等待轨道。"
        ],
        title="日志证据", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=13.8
    )
    add_image(slide, wait_diag, Inches(6.15), Inches(1.42), width=Inches(6.0))
    add_footer(slide, "推荐讲法：carrier 追的不是理想圆轨道，而是一个会因 TECS/loiter 波动而变化的真实轨迹。")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "4. 第三层根因：现在卡的是“末段相对速度”，不是“最小距离”", "之所以能进 DOCKING 却收不进 COMPLETED，是因为几何上靠近了，但动力学上还没有锁住")
    add_metric_card(slide, Inches(0.78), Inches(1.45), Inches(1.85), Inches(0.95), "best window", best["best_window_phase"],
                    accent=COLOR_BLUE, note=f"t = {best['best_window_t_sec']} s")
    add_metric_card(slide, Inches(2.85), Inches(1.45), Inches(1.85), Inches(0.95), "最小距离", f"{best['min_distance_m']} m",
                    accent=COLOR_ORANGE, note="已经靠近")
    add_metric_card(slide, Inches(4.92), Inches(1.45), Inches(2.05), Inches(0.95), "rel speed", f"{best['best_window_rel_speed_mps']} m/s",
                    accent=COLOR_RED, note="仍然偏大")
    add_bullets(
        slide, Inches(0.72), Inches(2.75), Inches(5.7), Inches(2.8),
        [
            "当前不是完全进不了 DOCKING，而是进了以后更像高速擦肩而过。",
            "DOCKING -> COMPLETED 需要同时满足小误差、小相对速度，还要保持若干控制周期。",
            "现在最近点到了，但 along-track 速度没压住，所以无法稳定进入真实捕获包络。"
        ],
        title="为什么只到 DOCKING", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.0
    )
    add_bullets(
        slide, Inches(0.72), Inches(5.8), Inches(5.7), Inches(0.8),
        [
            "所以真正的问题不是“距离还不够小”，而是“最近点附近速度还不够小”。"
        ],
        title="一句话", fill_color=RGBColor(232, 239, 251), line_color=COLOR_BLUE, bullet_size=14.0
    )
    add_bullets(
        slide, Inches(6.7), Inches(1.45), Inches(5.5), Inches(5.1),
        [
            "这也是为什么早期看起来“能盘旋、能 completed”的版本不能直接代表当前难度。",
            "现在这版更真实：fixed-wing 要先真实起飞、定高等待、只在末段减速配合。",
            "约束更强以后，问题从“链路通不通”变成“真实窗口能否稳定形成并被抓住”。",
            "所以当前主瓶颈已很清楚：等待轨道质量 + terminal speed matching。"
        ],
        title="为什么这是更难但也更有价值的问题", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=13.9
    )
    add_footer(slide, "结论：现在卡的是“真实飞控下的可捕获窗口”，不是“系统没接好”。")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    add_title(slide, "5. 所以下一步该改哪：按根因排序，不再盲调", "从“最影响结果的地方”往下做，而不是继续在外围参数上反复试错")
    add_bullets(
        slide, Inches(0.72), Inches(1.45), Inches(5.55), Inches(4.8),
        [
            "优先级 1：固定翼等待轨道稳定性。先让 30 m / 10 m/s 等待轨道更可预测。",
            "优先级 2：DOCKING 末段 along-track 速度匹配。先把最近点附近相对速度压下来。",
            "优先级 3：真正进入 terminal corridor 后，再让 fixed-wing 末段配合减速，而不是更早减速。",
            "优先级 4：继续把 TECS / loiter / glide 状态并入日志与控制判据，形成 TECS-aware 终端控制。"
        ],
        title="接下来最值得做的 4 件事", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.1
    )
    add_bullets(
        slide, Inches(6.7), Inches(1.45), Inches(5.55), Inches(4.8),
        [
            "不要再把时间主要花在“距离阈值再调一点”“start 时机再猜一下”上。",
            "现在最该做的是让目标 fixed-wing 本身变得更可预期，再让 terminal phase 变得更可捕获。",
            "如果这两件事做好，状态机和现有分层控制框架本身是有希望继续收进 0.2 m 的。"
        ],
        title="这轮定位后的工作原则", fill_color=COLOR_WHITE, line_color=RGBColor(220, 226, 235), bullet_size=14.1
    )
    add_textbox(
        slide, Inches(0.92), Inches(6.1), Inches(11.45), Inches(0.55),
        "推荐收束句：PX4 迁移没有失败，它已经把问题暴露到了更真实、更有研究价值的层面——真实 fixed-wing 飞控约束下的末段窗口稳定性。",
        size=15.5, bold=True, color=COLOR_NAVY, fill_color=RGBColor(232, 239, 251),
        line_color=COLOR_BLUE, align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE
    )
    add_footer(slide, "这页适合放在“问题分析”结尾，承接下一阶段计划。")

    return prs


def main() -> None:
    DOCS.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(OUTPUT_PPTX)
    prs.save(OUTPUT_PPTX_CN)
    print(OUTPUT_PPTX)
    print(OUTPUT_PPTX_CN)


if __name__ == "__main__":
    main()
